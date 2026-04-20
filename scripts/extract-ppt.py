#!/usr/bin/env python3
"""Extract text from .pptx files into structured markdown."""
import sys
from pathlib import Path
from pptx import Presentation


def extract_pptx(path: Path) -> str:
    """Return a markdown-formatted string with slide-by-slide text."""
    prs = Presentation(str(path))
    lines = [f"# {path.stem}\n"]
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"\n## Slide {idx}\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                for line in shape.text.splitlines():
                    line = line.strip()
                    if line:
                        lines.append(f"- {line}")
    return "\n".join(lines) + "\n"


def main() -> int:
    if len(sys.argv) != 3:
        print("Usage: extract-ppt.py <input.pptx> <output.md>", file=sys.stderr)
        return 1
    src = Path(sys.argv[1])
    dst = Path(sys.argv[2])
    if not src.exists():
        print(f"Not found: {src}", file=sys.stderr)
        return 2
    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text(extract_pptx(src), encoding="utf-8")
    print(f"Wrote {dst} ({dst.stat().st_size} bytes)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
